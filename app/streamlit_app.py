import sys
from pathlib import Path
PROJECT_ROOT = str(Path(__file__).resolve().parents[1])
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

# app/streamlit_app.py
import io
import os
from pathlib import Path
import streamlit as st
from dotenv import load_dotenv
from typing import List, Dict

load_dotenv()

# Ensure we run from project root so app.* imports resolve
# (Run streamlit from project root: `streamlit run app/streamlit_app.py`)
from app.rag import run_rag
# embed_index.query_index can be imported if you want query-only features
# from app.embed_index import query_index

# ---------- Utilities to create downloadable PDF / PPTX in-memory ----------
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation

def create_pdf_bytes(title: str, summary: str, sources: List[Dict]) -> bytes:
    """
    create a simple PDF in memory and return bytes
    """
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 740, f"Summary for: {title}")
    c.setFont("Helvetica", 11)
    text = c.beginText(72, 720)
    for line in summary.split("\n"):
        # naive wrap to 90 chars
        while len(line) > 90:
            text.textLine(line[:90])
            line = line[90:]
        text.textLine(line)
    c.drawText(text)
    c.showPage()
    c.setFont("Helvetica-Bold", 14)
    c.drawString(72, 740, "Sources")
    c.setFont("Helvetica", 10)
    y = 720
    for s in sources:
        line = f"[{s.get('idx')}] {s.get('title','')} — {s.get('url','')}"
        if y < 120:
            c.showPage(); y = 740
            c.setFont("Helvetica", 10)
        c.drawString(72, y, line[:240])
        y -= 14
    c.save()
    buf.seek(0)
    return buf.read()

def create_pptx_bytes(title: str, summary: str, sources: List[Dict]) -> bytes:
    """
    create a simple PPTX in memory and return bytes
    """
    buf = io.BytesIO()
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Summary: {title}"
    # summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    body = slide.shapes.placeholders[1].text_frame
    # add paragraphs (break by blank line)
    for para in summary.split("\n\n"):
        p = body.add_paragraph()
        p.text = para[:2000]
    # sources slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sources"
    body = slide.shapes.placeholders[1].text_frame
    for s in sources:
        p = body.add_paragraph()
        p.text = f"[{s.get('idx')}] {s.get('title','')} — {s.get('url','')}"
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ----------------- Streamlit UI -----------------
st.set_page_config(page_title="SmartNews Synthesizer — Streamlit", layout="wide")
st.title("SmartNews Synthesizer — Streamlit UI")
st.markdown("Enter a query, fetch news from NewsAPI, run retrieval, and generate a verified news-style summary using Ollama.")

# Sidebar controls
with st.sidebar:
    st.header("Options")
    days = st.number_input("Search last N days", min_value=1, max_value=30, value=1)
    top_k = st.number_input("Top-k passages to retrieve", min_value=1, max_value=10, value=int(os.getenv("RAG_TOP_K", 5)))
    max_tokens = st.number_input("Ollama max tokens", min_value=64, max_value=2000, value=int(os.getenv("OLLAMA_MAX_TOKENS", 512)))
    run_button = st.button("Run RAG & Generate")

# Main input
query = st.text_input("Query", placeholder="e.g. 'AI regulation Europe'")

# space for outputs
summary_area = st.empty()
sources_area = st.container()
evidence_area = st.container()
download_area = st.container()
log_area = st.empty()

def pretty_show_sources(sources):
    for s in sources:
        st.write(f"**[{s.get('idx')}]** {s.get('title')} — {s.get('url')}")

if run_button:
    if not query.strip():
        st.warning("Please enter a query before running.")
    else:
        with st.spinner("Running retrieval and calling Ollama... (this can take a few seconds)"):
            try:
                # run_rag returns dict with summary_text, sources, retrieved
                result = run_rag(query, top_k=top_k, model_name=os.getenv("OLLAMA_MODEL"), max_tokens=max_tokens, verbose=False)
            except Exception as e:
                st.error(f"Error during RAG: {e}")
                st.stop()

        # display summary
        st.subheader("Generated Summary")
        st.write(result.get("summary_text", "No text returned."))

        # show claims / raw output if available
        st.subheader("Raw Ollama Output (first 2000 chars)")
        st.code(str(result.get("raw_ollama"))[:2000])

        # show sources mapping
        st.subheader("Sources (citation mapping)")
        pretty_show_sources(result.get("sources", []))

        # show retrieved passages (expanders)
        st.subheader("Retrieved Passages (evidence)")
        retrieved = result.get("retrieved", [])
        for r in retrieved:
            md = r.get("metadata", {})
            header = f"[{md.get('doc_id')[:8]}] {md.get('title','Untitled')} — {md.get('source','')}"
            with st.expander(header):
                st.write(f"Chunk index: {md.get('chunk_index')}")
                st.write(md.get("chunk_text", "(no chunk text)"))
                st.caption(f"Score: {r.get('score')}")

        # create downloadable files
        pdf_bytes = create_pdf_bytes(query, result.get("summary_text",""), result.get("sources", []))
        pptx_bytes = create_pptx_bytes(query, result.get("summary_text",""), result.get("sources", []))

        st.subheader("Download")
        col1, col2 = st.columns(2)
        with col1:
            st.download_button("Download PDF", data=pdf_bytes, file_name="summary.pdf", mime="application/pdf")
        with col2:
            st.download_button("Download PPTX", data=pptx_bytes, file_name="summary.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        st.success("Done — check Summary, Sources, and Evidence above.")

# Help / troubleshooting
with st.expander("Troubleshooting & notes"):
    st.write("""
    * Run this app from the project root so that `app.*` imports resolve:
      `streamlit run app/streamlit_app.py`
    * Make sure your venv is activated and Ollama is running (`ollama serve`) and `llama3:latest` is installed.
    * If the summary seems to invent facts, reduce top_k or make the prompt stricter in `app/rag.py`.
    * The app uses the `run_rag` function you created earlier — any changes you make to that function will reflect here.
    """)

